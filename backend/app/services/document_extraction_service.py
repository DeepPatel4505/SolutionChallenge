import io

import httpx

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


async def _fetch_bytes(file_url: str) -> bytes:
    async with httpx.AsyncClient(timeout=120.0, follow_redirects=True) as client:
        resp = await client.get(file_url)
        if resp.status_code != 200:
            raise RuntimeError(f"Failed to download file: {resp.status_code} - {resp.text}")
        return resp.content


def _extract_pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts).strip()


def _extract_docx_text(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    parts: list[str] = []

    # Paragraphs
    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if text:
            parts.append(text)

    # Tables (optional but helps content-heavy DOCX)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join((cell.text or "").strip() for cell in row.cells)
            row_text = row_text.strip()
            if row_text:
                parts.append(row_text)

    return "\n".join(parts).strip()


def _extract_pptx_text(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # Many PPT text containers have `.text`
            text = ""
            if hasattr(shape, "text"):
                try:
                    text = (shape.text or "").strip()
                except Exception:
                    text = ""
            if text:
                parts.append(text)

        # Add speaker notes too (if present)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = (notes_frame.text or "").strip()
            if notes_text:
                parts.append(notes_text)

    return "\n\n".join(parts).strip()


async def extract_document_text(file_url: str, file_ext: str) -> dict:
    """
    Extract text from a document URL.
    - PDF: pypdf
    - DOCX: python-docx
    - PPTX: python-pptx

    No OCR is performed. If the document has no selectable text, extraction will be empty.
    """
    ext = (file_ext or "").lower().lstrip(".")
    data = await _fetch_bytes(file_url)

    if ext == "pdf":
        text = _extract_pdf_text(data)
    elif ext == "docx":
        text = _extract_docx_text(data)
    elif ext == "pptx":
        text = _extract_pptx_text(data)
    else:
        raise RuntimeError(f"Unsupported document extension: {ext}")

    if not text:
        raise RuntimeError("No extractable text found in the uploaded document (OCR not enabled).")

    return {"transcript_text": text, "transcript_json": None}

